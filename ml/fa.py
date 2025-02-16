from fastapi import FastAPI, UploadFile, File
from pydantic import BaseModel
from typing import List, Dict
from openai import OpenAI
import json
import faiss
import os
import numpy as np
from utils import extract_json, convert_width_height_to_numbers, process_json, translate_to_english, chunk_text
from PIL import Image
import uuid
from diffusers import AutoPipelineForText2Image
from sentence_transformers import SentenceTransformer
import requests
import logging
from fastapi.middleware.cors import CORSMiddleware
import torch
from transformers import pipeline
import PyPDF2
from pptx import Presentation
from PIL import Image
import pytesseract
import io

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  
    allow_credentials=True,
    allow_methods=["*"], 
    allow_headers=["*"], 
)

embedding_dim = 312
threshold = 0.7
d = 312

index_file = "faiss.index"
SAVE_DIR = os.path.abspath("generated_images")  
os.makedirs(SAVE_DIR, exist_ok=True)

image_index_file = "faiss_image.index"
image_paths_file = "image_paths.json"
image_index = faiss.IndexFlatL2(d)
image_paths = []
IMG_DIR = os.path.abspath("img")
os.makedirs(IMG_DIR, exist_ok=True)

try:
    image_embeddings = []
    for img_name in os.listdir(IMG_DIR):
        img_path = os.path.join(IMG_DIR, img_name)
        embedding = encoder.encode(img_name).tolist()
        image_embeddings.append(embedding)
        image_paths.append(img_path)
    if image_embeddings:
        image_index.add(np.array(image_embeddings, dtype=np.float32))
        faiss.write_index(image_index, image_index_file)
    with open(image_paths_file, "w", encoding="utf-8") as f:
        json.dump(image_paths, f, ensure_ascii=False, indent=2)
    logger.info(f"Загружено {len(image_paths)} изображений в индекс.")
except Exception as e:
    logger.error(f"Ошибка при загрузке изображений в индекс: {e}")


encoder = SentenceTransformer("cointegrated/rubert-tiny2")
pipe = AutoPipelineForText2Image.from_pretrained("kandinsky-community/kandinsky-2-1", torch_dtype=torch.float16)
pipe.enable_model_cpu_offload()
summarizer = pipeline("summarization", model="d0rj/rut5-base-summ")

if os.path.exists(index_file):
    index = faiss.read_index(index_file)
else:
    index = faiss.IndexFlatL2(d)
try:
    with open("slide_contents.json", "r", encoding="utf-8") as f:
        slide_contents = json.load(f)
except:
    slide_contents = []

class SlideData(BaseModel):
    id: int
    title: str
    content: str
    embedding: List[float]

class SearchRequest(BaseModel):
    query_embedding: List[float]
    top_k: int

class SlideRequest(BaseModel):
    topic: str
    description: str

class LayoutRequest(BaseModel):
    slides: List[Dict]

class GenerateRequest(BaseModel):
    prompt: str
    negative_prompt: str = "low quality, bad quality"
    height: int = 256
    width: int = 256
    prior_guidance_scale: float = 1.0

async def generate_presentation_structure(topic: str, description: str, top_k=1):
    logger.info("Векторизация темы")
    topic_emb = encoder.encode(topic).tolist()
    logger.info("Поиск дополнительного контента")
    search_results = await search_slides(SearchRequest(query_embedding=topic_emb, top_k=top_k))
    logger.info(f"Найдено: {search_results}")
    messages = [
        {
            "role": "user",
            "content": f"Создай структуру презентации на тему '{topic}'. Описание: {description}. Дополнительная релевантная информация: {search_results}. Проверь cогласованность слайдов, наличие начального и конечного слайдов. Проверь, что весь текст на русском языке. Проверь формат ответа и корректность json, чтобы его можно было выделить по скобкам {}.\\n" + 
                       """Ответ в JSON-формате: 
                       {
                          "slides": [
                            {
                              "elements": [
                                {
                                  "type": "title",
                                  "content": "Наименование",
                                },
                                {
                                  "type": "text",
                                  "content": "Текст",

                                },
                                {
                                  "type": "list",
                                  "content": "Содержимое списка через",

                                },
                                {
                                  "type": "image",
                                  "content": "Тематика изображения",
                                }
                              ]
                            },
                            ...,
                            ...
                          ]
                        }
                        """
        }
    ]
    response = requests.post("http://localhost:8088/v1/completions", json={
        "model": "Qwen/Qwen2.5-7B-Instruct",
        "prompt": messages[0]["content"], 
        "max_tokens": 4096})
    
    response_text = response.json()["choices"][0]["text"]
    return extract_json(response_text)

def enrich_slide_layout(slides: List[Dict]):
    messages = [
        {
            "role": "user",
            "content": f"Добавь размеры (разметку) (поля width и heights в процентах) для каждого элемента слайда в JSON: {json.dumps(slides, ensure_ascii=False)}. Проверь, что элементы расположены корректно. Верни ответ в формате JSON."
        }
    ]
    response = requests.post("http://localhost:8088/v1/completions", json={
        "model": "Qwen/Qwen2.5-7B-Instruct",
        "prompt": messages[0]["content"], 
        "max_tokens": 4096})
    
    response_text = response.json()["choices"][0]["text"]
    return response_text

def process_pdf(file: UploadFile):
    pdf_reader = PyPDF2.PdfReader(file.file)
    full_text = " "
    for page in pdf_reader.pages:
        full_text += page.extract_text() + " "
    
    chunks = chunk_text(full_text)
    summarized_chunks = [summarizer(chunk, max_length=200, min_length=50, do_sample=False)[0]['summary_text'] for chunk in chunks]
    
    for chunk in summarized_chunks:
        embedding = encoder.encode(chunk).tolist()
        slide = SlideData(id=len(slide_contents), title="PDF Summary", content=chunk, embedding=embedding)
        index.add(np.array([slide.embedding], dtype=np.float32))
        slide_contents.append(slide.content)
        with open("slide_contents.json", "w", encoding="utf-8") as f:
            json.dump(slide_contents, f, ensure_ascii=False, indent=2)
        faiss.write_index(index, index_file)
    
    return {"chunks": summarized_chunks}


def process_pptx(file: UploadFile):
    with io.BytesIO(file.file.read()) as file_stream:
        prs = Presentation(file_stream)
        slides_data = []
        for slide in prs.slides:
            title = ""
            content = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if shape.text_frame and shape.text_frame.paragraphs:
                        if shape.text_frame.paragraphs[0].level == 0:
                            title = text
                        else:
                            content += " " + text
                if shape.shape_type == 13:  
                    img_data = shape.image.blob
                    with io.BytesIO(img_data) as img_io:
                        with Image.open(img_io) as img:
                            ocr_text = pytesseract.image_to_string(img)
                            if ocr_text.strip():
                                content += " " + ocr_text.strip()
            slides_data.append({"title": title, "content": content.strip()})
        for slide in slides_data:
            if len(slide["content"]) > 50:
                embedding = encoder.encode(slide["content"]).tolist()
                slide_data = SlideData(id=len(slide_contents), title=slide["title"], content=slide["content"], embedding=embedding)
                index.add(np.array([slide_data.embedding], dtype=np.float32))
                slide_contents.append(slide_data.content)
                with open("slide_contents.json", "w", encoding="utf-8") as f:
                    json.dump(slide_contents, f, ensure_ascii=False, indent=2)
                faiss.write_index(index, index_file)
        return {"slides": slides_data}

async def find_similar_image(prompt: str):
    embedding = encoder.encode(prompt).astype(np.float32).reshape(1, -1)
    distances, indices = image_index.search(embedding, 1)
    if distances[0][0] < threshold and indices[0][0] < len(image_paths):
        return image_paths[indices[0][0]]
    return None

@app.post("/enrich_layout/")
async def enrich_layout(request: LayoutRequest):
    enriched_layout = enrich_slide_layout(request.slides)
    extracted_json = extract_json(enriched_layout)
    return {"slides": extracted_json} if "slides" not in extracted_json else extracted_json

@app.post("/generate_presentation/")
async def generate_presentation(request: SlideRequest):
    presentation = await generate_presentation_structure(request.topic, request.description)
    return presentation

@app.post("/index_slide/")
async def index_slide(slide: SlideData):
    index.add(np.array([slide.embedding], dtype=np.float32))
    faiss.write_index(index, index_file)
    slide_contents.append(slide.content)
    with open(slide_contents_file, "w", encoding="utf-8") as f:
        json.dump(slide_contents, f, ensure_ascii=False, indent=2)
    return {"message": "Slide indexed successfully"}

@app.post("/search_slides/")
async def search_slides(request: SearchRequest):
    query_vector = np.array([request.query_embedding], dtype=np.float32)
    distances, indices = index.search(query_vector, request.top_k)
    # retrieved_texts = [slide_contents[idx] for idx in indices.flatten() if idx < len(slide_contents)]ё
    retrieved_texts = []
    for dist, idx in zip(distances.flatten(), indices.flatten()):
        if idx < len(slide_contents) and dist <= threshold:
            retrieved_texts.append(slide_contents[idx])
    return retrieved_texts

@app.post("/main/")
async def process_presentation(request: SlideRequest):
    presentation_response = await generate_presentation(request)
    enrich_response = await enrich_layout(LayoutRequest(slides=presentation_response["slides"]))
    if enrich_response:
        converted_data = convert_width_height_to_numbers(enrich_response)
        converted_data = process_json(converted_data)
        for slide in converted_data["slides"]:
            for element in slide["elements"]:
                if element["type"] == "image":
                    prompt = element["content"]
                    logger.info(f"prompt: {prompt}")
                    existing_image = await find_similar_image(prompt)
                    if existing_image:
                        logger.info(f"Найдено похожее изображение: {existing_image}")
                        element["content"] = existing_image
                    else:
                        prompt = translate_to_english(prompt)
                        logger.info(f"prompt: {prompt}")
                        image_path = await generate_image(GenerateRequest(prompt=prompt))
                        element["content"] = image_path["image_path"] 
        return converted_data
    else:
        return {"error": "Не удалось обработать JSON"}


@app.post("/process_file/")
async def process_file(file: UploadFile = File(...)):
    file_ext = file.filename.split(".")[-1].lower()
    
    if file_ext == "pptx":
        return process_pptx(file)
    elif file_ext == "pdf":
        return process_pdf(file)
    else:
        return {"error": "Unsupported file format"}

@app.post("/generate_image/")
async def generate_image(request: GenerateRequest):
    image = pipe(
        prompt=request.prompt,
        negative_prompt=request.negative_prompt,
        prior_guidance_scale=request.prior_guidance_scale,
        height=request.height,
        width=request.width
    ).images[0]
    
    file_name = f"{uuid.uuid4().hex}.png"
    file_path = os.path.join(SAVE_DIR, file_name)
    image.save(file_path)
    
    logger.info(f"Сгенерировано изображение: {file_path}")
    return {"image_path": file_path}


@app.post("/clear_index/")
async def clear_index():
    slide_contents = []
    with open(slide_contents_file, "w", encoding="utf-8") as f:
        json.dump(slide_contents, f, ensure_ascii=False, indent=2)
    index = faiss.IndexFlatL2(embedding_dim)
    faiss.write_index(index, "faiss.index")
    return {"Info": "cleared"}
